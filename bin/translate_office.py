import os
import subprocess
import openpyxl
import docx
import fitz  # PyMuPDF
from pptx import Presentation
from deep_translator import GoogleTranslator

def translate_text(text, translator):
    """Helper function to translate text while handling empty strings."""
    if not text or not isinstance(text, str) or not text.strip():
        return text
    try:
        return translator.translate(text)
    except Exception as e:
        print(f"  [!] Failed to translate text: '{text[:30]}...'. Error: {e}")
        return text

def process_excel(file_path, output_path, translator):
    print("Processing Excel workbook...")
    wb = openpyxl.load_workbook(file_path)
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        for row in ws.iter_rows():
            for cell in row:
                if cell.value and isinstance(cell.value, str):
                    cell.value = translate_text(cell.value, translator)
    wb.save(output_path)

def process_word(file_path, output_path, translator):
    print("Processing Word document...")
    doc = docx.Document(file_path)

    # Translate standard paragraphs
    for para in doc.paragraphs:
        if para.text.strip():
            para.text = translate_text(para.text, translator)

    # Translate text inside tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    if para.text.strip():
                        para.text = translate_text(para.text, translator)

    doc.save(output_path)

def process_powerpoint(file_path, output_path, translator):
    print("Processing PowerPoint presentation...")
    prs = Presentation(file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.text.strip():
                    paragraph.text = translate_text(paragraph.text, translator)

    prs.save(output_path)

def process_pdf(file_path, output_path, translator):
    print("Extracting and translating PDF (saving to Word format)...")
    pdf_doc = fitz.open(file_path)
    new_doc = docx.Document()

    for page_num in range(len(pdf_doc)):
        print(f"  Processing page {page_num + 1}...")
        page = pdf_doc.load_page(page_num)
        text = page.get_text("text")

        # Split text into manageable paragraphs to maintain some structure
        paragraphs = text.split('\n\n')

        for para in paragraphs:
            clean_para = para.strip().replace('\n', ' ')
            if clean_para:
                translated = translate_text(clean_para, translator)
                new_doc.add_paragraph(translated)

        # Add a page break between PDF pages
        if page_num < len(pdf_doc) - 1:
            new_doc.add_page_break()

    pdf_doc.close()
    new_doc.save(output_path)

def main():
    print("--- The Ultimate macOS Office & PDF Translator ---")

    file_path = input("Enter the path to the file (.xlsx, .docx, .doc, .pptx, .pdf): ").strip().strip("'\"")

    if not os.path.exists(file_path):
        print("Error: File not found.")
        return

    src_lang = input("Enter source language code (or 'auto'): ").strip().lower()
    tgt_lang = input("Enter target language code (e.g., 'en'): ").strip().lower()

    try:
        translator = GoogleTranslator(source=src_lang, target=tgt_lang)
    except Exception as e:
        print(f"Error initializing translator: {e}")
        return

    # Determine file type
    dir_name, file_name = os.path.split(file_path)
    name, ext = os.path.splitext(file_name)
    ext = ext.lower()

    # Handle output extensions (force PDF and old DOC into modern DOCX)
    if ext == '.pdf' or ext == '.doc':
        out_ext = '.docx'
    else:
        out_ext = ext

    output_filename = f"{name}_{tgt_lang}{out_ext}"
    output_path = os.path.join(dir_name, output_filename)

    print(f"\nLoading {file_name}...")

    # Route to the correct function based on the file extension
    try:
        if ext == '.xlsx':
            process_excel(file_path, output_path, translator)
        elif ext == '.docx':
            process_word(file_path, output_path, translator)
        elif ext == '.doc':
            print("  Legacy .doc format detected. Upgrading to .docx using macOS textutil...")
            # Use macOS built-in textutil to convert the file
            subprocess.run(['textutil', '-convert', 'docx', file_path], check=True)

            # Find the newly created .docx file (same name, new extension)
            temp_docx_path = os.path.join(dir_name, name + ".docx")

            # Process the new modern file
            process_word(temp_docx_path, output_path, translator)

            # Delete the temporary un-translated modern file to keep the folder clean
            if os.path.exists(temp_docx_path):
                os.remove(temp_docx_path)

        elif ext == '.pptx':
            process_powerpoint(file_path, output_path, translator)
        elif ext == '.pdf':
            process_pdf(file_path, output_path, translator)
        else:
            print(f"Error: Unsupported file format '{ext}'.")
            return

        print(f"\nSuccess! Translated file saved to:\n{output_path}")

    except Exception as e:
        print(f"\nAn error occurred while processing the file: {e}")

if __name__ == "__main__":
    main()
